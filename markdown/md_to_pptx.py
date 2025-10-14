#!/usr/bin/env python3
"""
md_to_pptx.py

Convert a markdown deck into a PowerPoint (.pptx).
- Slides are separated by lines containing only `---`
- The first '#' heading on each slide becomes the slide title
- Bullets are lines starting with -, *, or numbered list markers
- Indentation controls bullet levels (2 spaces = level 1, 4 spaces = level 2, etc.)
- Fenced code blocks (``` ... ```) are inserted as monospace paragraphs

Usage:
    python md_to_pptx.py path/to/deck.md
Outputs:
    path/to/deck.pptx
"""

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# ---------- Parsing helpers ----------

HEADING_RE = re.compile(r'^\s*#{1,6}\s+(.*)$')
BULLET_RE  = re.compile(r'^\s*([-*]|(\d+[\.\)]))\s+(.*)$')
HRULE_RE   = re.compile(r'^\s*---\s*$')  # slide separator
FRONTMATTER_DELIM = re.compile(r'^\s*---\s*$')  # YAML front matter delimiter

def strip_front_matter(lines):
    """Remove leading YAML front matter block delimited by '---'."""
    if not lines:
        return lines
    if FRONTMATTER_DELIM.match(lines[0]):
        # find next delimiter
        for i in range(1, len(lines)):
            if FRONTMATTER_DELIM.match(lines[i]):
                return lines[i+1:]
    return lines

def split_slides(md_text):
    """Split markdown text into slide chunks using lines containing only '---'."""
    # Normalize newlines
    lines = md_text.replace('\r\n', '\n').replace('\r', '\n').split('\n')

    # Strip optional YAML front matter at top
    lines = strip_front_matter(lines)

    slides = []
    current = []
    for ln in lines:
        if HRULE_RE.match(ln):
            # separator: commit current slide
            slides.append('\n'.join(current).strip())
            current = []
        else:
            current.append(ln)
    # last slide
    if current:
        slides.append('\n'.join(current).strip())

    # Remove empties
    return [s for s in slides if s.strip()]

def parse_slide(chunk):
    """
    Parse a slide chunk into (title, content_items)
    content_items is a list of dicts:
      { "type": "bullet", "text": str, "level": int, "mono": bool }
      { "type": "code",   "text": str }  # fenced block
    """
    lines = chunk.splitlines()
    title = None
    items = []
    in_code = False
    code_fence = None
    code_buffer = []

    def flush_code():
        nonlocal code_buffer, items
        if code_buffer:
            items.append({"type": "code", "text": '\n'.join(code_buffer)})
            code_buffer = []

    for raw in lines:
        ln = raw.rstrip('\n')

        # fenced code block start/end
        if not in_code and ln.strip().startswith("```"):
            in_code = True
            code_fence = ln.strip()
            code_buffer = []
            continue
        elif in_code and ln.strip().startswith("```"):
            in_code = False
            flush_code()
            code_fence = None
            continue

        if in_code:
            code_buffer.append(ln)
            continue

        # detect first heading as title
        if title is None:
            m = HEADING_RE.match(ln)
            if m:
                title = m.group(1).strip()
                continue

        # bullets ( -, *, 1. 1) ) with indentation -> level
        m = BULLET_RE.match(ln)
        if m:
            # Compute level by leading spaces (2 spaces per level)
            leading_spaces = len(ln) - len(ln.lstrip(' '))
            level = min((leading_spaces // 2), 5)
            text = m.group(3).strip()
            items.append({"type": "bullet", "text": text, "level": level, "mono": False})
            continue

        # Additional text lines (non-empty) → treat as bullets at level 0
        if ln.strip():
            items.append({"type": "bullet", "text": ln.strip(), "level": 0, "mono": False})

    # If no title, use the first non-empty line
    if title is None:
        for ln in lines:
            if ln.strip():
                title = ln.strip()
                break
        if title is None:
            title = "Untitled"

    return title, items

# ---------- PPTX helpers ----------

def add_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()  # remove default paragraph

    # We put bullets in the main text frame; code blocks as monospace paragraphs too
    first_para_added = False
    for it in items:
        if it["type"] == "bullet":
            if not first_para_added:
                p = tf.paragraphs[0]
                first_para_added = True
            else:
                p = tf.add_paragraph()
            p.text = it["text"]
            p.level = max(0, min(5, it.get("level", 0)))
            # Optional: tweak font size slightly
            for r in p.runs:
                r.font.size = Pt(18)
        elif it["type"] == "code":
            # Add each code block as a paragraph with monospace font
            p = tf.add_paragraph()
            p.text = it["text"]
            p.level = 0
            # monospace styling
            for r in p.runs:
                r.font.name = "Courier New"
                r.font.size = Pt(14)

def md_to_pptx(md_path: Path):
    if not md_path.exists() or not md_path.is_file():
        raise FileNotFoundError(f"Markdown file not found: {md_path}")

    out_path = md_path.with_suffix(".pptx")

    md_text = md_path.read_text(encoding="utf-8")
    chunks = split_slides(md_text)
    slides = [parse_slide(c) for c in chunks]

    prs = Presentation()
    # Remove the default first slide if any templates insert one (we're adding our own)
    # (python-pptx starts empty with standard template; no need to delete)

    for title, items in slides:
        add_slide(prs, title, items)

    prs.save(out_path)
    print(f"Saved: {out_path}")

def main():
    ap = argparse.ArgumentParser(description="Convert Markdown slides (--- separators) to PPTX.")
    ap.add_argument("markdown_file", help="Path to the input .md file")
    args = ap.parse_args()

    md_path = Path(args.markdown_file).expanduser().resolve()
    md_to_pptx(md_path)

if __name__ == "__main__":
    main()
